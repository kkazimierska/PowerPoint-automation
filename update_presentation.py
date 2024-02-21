from pptx import Presentation
from datetime import datetime
import settings
from update_facebook import get_data
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pptx
from duplicate import duplicate_slide


# TODO: Task1: Duplicate slide
# TODO: Task1: Iterate over data rows


## Get data
team_competency = get_data(settings.INPUT_EXCEL)

## Open pptx and explore
prs = Presentation('Facebookv1.pptx')

## Replace the picture from shape = 1 to "images/dev.jpg"
img_path = "images/dev.jpg"
img2 = pptx.parts.image.Image.from_file(img_path)


## CREATE X slides

n = len(team_competency.index)
prs = duplicate_slide(pres=prs, index=5)
# Iterate change image and adding the data
slide5 = prs.slides[5]

# what image you're actually changing...
img_shape = slide5.shapes[1]

def change_image(old_image_shape, new_image):
    """Change the image for the shape image

    :param old_image_shape: Default template image
    :param new_image: Desired image
    :return: _description_
    """
    # get part and rId from shape we need to change
    slide_part, rId = old_image_shape.part, old_image_shape._element.blip_rId
    image_part = slide_part.related_part(rId)

    # overwrite old blob info with new blob info
    image_part.blob = new_image._blob
    return old_image_shape

img_shape = change_image(old_image_shape=img_shape, new_image=img2)


def set_data_to_slide(slide, team_competency, row_num: int):
    """Set the data from csv row to pptx slide.

    :param slide: _description_
    :param team_competency: _description_
    """
    slide.shapes[6].text = str(team_competency['Education'][row_num])
    slide.shapes[7].text = str(team_competency['Professional Interest'][row_num])
    slide.shapes[8].text = str(team_competency['Delivery Org 1'][row_num])
    slide.shapes[9].text = str(team_competency['Position'][row_num])
    slide.shapes[9].text_frame.paragraphs[0].font.bold = True
    slide.shapes[10].text = str(team_competency['Experience'][row_num])
    slide.shapes[11].text = str(team_competency['_professional_competencies'][row_num])
    slide.shapes[12].text = str(f"{team_competency['Name'][row_num]} ({team_competency['Initials'][row_num]})")
    slide.shapes[12].text_frame.paragraphs[0].font.bold = True
    return slide

slide5 = set_data_to_slide(slide=slide5, team_competency=team_competency, row_num=1)

ts = datetime.now()

prs.save(f"Draft_readed{ts.microsecond}.pptx")